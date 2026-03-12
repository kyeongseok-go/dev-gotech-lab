#!/usr/bin/env python3
"""gotech.lab 1주 MVP 스프린트 요약 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- 색상 팔레트 ---
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # 짙은 남색 배경
BG_CARD = RGBColor(0x1A, 0x24, 0x3B)       # 카드 배경
ACCENT = RGBColor(0x60, 0xA5, 0xFA)        # 파란 강조
ACCENT2 = RGBColor(0x34, 0xD3, 0x99)       # 초록 강조
ACCENT3 = RGBColor(0xF4, 0x72, 0xB6)       # 핑크 강조
ACCENT4 = RGBColor(0xFB, 0xBF, 0x24)       # 노란 강조
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY = RGBColor(0x47, 0x55, 0x69)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=BG_DARK):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, radius=None):
    """둥근 사각형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_para(text_frame, text, size=14, color=LIGHT_GRAY, bold=False, space_before=Pt(4), bullet=False):
    """텍스트 프레임에 단락 추가"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "맑은 고딕"
    p.space_before = space_before
    if bullet:
        p.level = 0
    return p


def add_accent_line(slide, left, top, width, color=ACCENT):
    """강조 라인 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_badge(slide, left, top, text, bg_color=ACCENT, text_color=BG_DARK, width=None):
    """뱃지 추가"""
    w = width or Inches(1.6)
    shape = add_shape(slide, left, top, w, Inches(0.35), bg_color, radius=0.5)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "맑은 고딕"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


# ============================================================
# 슬라이드 1: 표지
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)

# 장식 원
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(8), Inches(8))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
circle.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(3), Inches(7), Inches(7))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x16, 0x20, 0x33)
circle2.line.fill.background()

add_accent_line(slide, Inches(1.5), Inches(2.3), Inches(1.2), ACCENT)

add_text(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
         "gotech.lab", size=54, color=WHITE, bold=True)
add_text(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
         "1주 MVP 스프린트 결과 보고서", size=28, color=ACCENT)
add_text(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.6),
         "개인 홈페이지 · 기술 블로그 · 서비스 플랫폼", size=18, color=MID_GRAY)
add_text(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.5),
         "2026.03 | Claude Code + Next.js + Cloudflare", size=14, color=DARK_GRAY)


# ============================================================
# 슬라이드 2: 프로젝트 개요
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(0.8), ACCENT)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
         "프로젝트 개요", size=36, color=WHITE, bold=True)

# 왼쪽: 목적
card = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), BG_CARD, 0.02)
add_text(slide, Inches(1.2), Inches(2.0), Inches(4.8), Inches(0.5),
         "🎯  프로젝트 목적", size=20, color=ACCENT, bold=True)

items = [
    "개인 기술 블로그 + 포트폴리오 사이트 구축",
    "AI(Claude Code)를 페어 프로그래머로 활용한 1주 스프린트",
    "Cloudflare 풀스택 (Pages + Workers + D1 + R2)",
    "채용 담당자 · 기술 리더 · 개발자 커뮤니티 대상",
]
txBox = add_text(slide, Inches(1.2), Inches(2.6), Inches(4.8), Inches(3.5), "", size=14)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in items:
    add_para(tf, f"▸  {item}", size=14, color=LIGHT_GRAY, space_before=Pt(10))

# 오른쪽: 핵심 수치
card2 = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.8), BG_CARD, 0.02)
add_text(slide, Inches(7.2), Inches(2.0), Inches(4.8), Inches(0.5),
         "📊  핵심 수치", size=20, color=ACCENT2, bold=True)

metrics = [
    ("7일", "스프린트 기간"),
    ("12개", "페이지 라우트"),
    ("8개", "MDX 콘텐츠"),
    ("3개", "D1 테이블"),
    ("20+", "수정 파일/일"),
]
y_start = 2.7
for i, (num, label) in enumerate(metrics):
    add_text(slide, Inches(7.4), Inches(y_start + i * 0.7), Inches(1.5), Inches(0.5),
             num, size=28, color=ACCENT4, bold=True)
    add_text(slide, Inches(9.0), Inches(y_start + i * 0.7 + 0.05), Inches(3), Inches(0.5),
             label, size=16, color=LIGHT_GRAY)


# ============================================================
# 슬라이드 3: 기술 스택
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(0.8), ACCENT2)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
         "기술 스택", size=36, color=WHITE, bold=True)

stacks = [
    ("프레임워크", ACCENT, [
        "Next.js 16 (App Router)",
        "TypeScript",
        "React 19",
    ]),
    ("스타일 / UI", ACCENT2, [
        "Tailwind CSS v4",
        "shadcn/ui (radix-nova)",
        "CSS Variables (oklch)",
    ]),
    ("콘텐츠", ACCENT3, [
        "MDX + Velite",
        "rehype-pretty-code",
        "Shiki (코드 하이라이팅)",
    ]),
    ("인프라 / 배포", ACCENT4, [
        "Cloudflare Pages + Workers",
        "D1 (SQLite) + R2 (Storage)",
        "OpenNext Adapter",
    ]),
]

for i, (title, color, items) in enumerate(stacks):
    x = 0.8 + i * 3.1
    card = add_shape(slide, Inches(x), Inches(1.8), Inches(2.8), Inches(4.5), BG_CARD, 0.02)
    add_accent_line(slide, Inches(x + 0.3), Inches(2.0), Inches(0.6), color)
    add_text(slide, Inches(x + 0.3), Inches(2.2), Inches(2.2), Inches(0.5),
             title, size=18, color=color, bold=True)
    txBox = add_text(slide, Inches(x + 0.3), Inches(2.8), Inches(2.2), Inches(3), "", size=13)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for item in items:
        add_para(tf, f"•  {item}", size=13, color=LIGHT_GRAY, space_before=Pt(8))


# ============================================================
# 슬라이드 4: Day 0~2 타임라인
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(0.8), ACCENT)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
         "스프린트 타임라인 — Day 0 ~ 2", size=36, color=WHITE, bold=True)

days_1 = [
    ("Day 0", "초기화", ACCENT, [
        "GitHub 저장소 생성",
        "Claude Code + 규칙 파일 설정",
        "Cloudflare 인증 초기화",
    ]),
    ("Day 1", "프로젝트 구성", ACCENT2, [
        "Next.js + Cloudflare + Tailwind + shadcn 초기화",
        "SiteHeader / SiteFooter / PageContainer 생성",
        "6개 페이지 라우트 스텁",
        "pnpm preview 기동 확인",
    ]),
    ("Day 2", "콘텐츠 파이프라인", ACCENT3, [
        "Velite 설치 + 3개 컬렉션 스키마",
        "MDXContent 렌더링 컴포넌트",
        "블로그/프로젝트/Showcase SSG",
        "홈 최신 콘텐츠 섹션",
    ]),
]

for i, (day, title, color, items) in enumerate(days_1):
    x = 0.8 + i * 4.0
    card = add_shape(slide, Inches(x), Inches(1.8), Inches(3.7), Inches(4.8), BG_CARD, 0.02)
    add_badge(slide, Inches(x + 0.3), Inches(2.0), day, color, BG_DARK, Inches(1.2))
    add_text(slide, Inches(x + 0.3), Inches(2.5), Inches(3.1), Inches(0.5),
             title, size=18, color=WHITE, bold=True)
    txBox = add_text(slide, Inches(x + 0.3), Inches(3.1), Inches(3.1), Inches(3), "", size=12)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for item in items:
        add_para(tf, f"▸  {item}", size=12, color=LIGHT_GRAY, space_before=Pt(6))


# ============================================================
# 슬라이드 5: Day 3~4 타임라인
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(0.8), ACCENT)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
         "스프린트 타임라인 — Day 3 ~ 4", size=36, color=WHITE, bold=True)

days_2 = [
    ("Day 3", "블로그 고도화", ACCENT4, [
        "rehype-pretty-code + Shiki 코드 하이라이팅",
        "Callout (info/warning/tip) 컴포넌트",
        "Table 스타일링 + Image figure 래퍼",
        "읽기 시간 + 이전/다음 글 네비게이션",
        "h2/h3 기반 정적 TOC",
        "카테고리·태그 필터 (searchParams)",
    ]),
    ("Day 4", "페이지 확장", ACCENT, [
        "/about: 경력 타임라인, 기술 스택, 관심 분야",
        "SubscribeForm 재사용 컴포넌트",
        "/subscribe 독립 페이지",
        "/showcase/[slug] 상세 (SSG, 뱃지, 링크)",
        "프로젝트 상세 보강 (techStack, repo/demo)",
    ]),
]

for i, (day, title, color, items) in enumerate(days_2):
    x = 0.8 + i * 6.0
    card = add_shape(slide, Inches(x), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD, 0.02)
    add_badge(slide, Inches(x + 0.3), Inches(2.0), day, color, BG_DARK, Inches(1.2))
    add_text(slide, Inches(x + 0.3), Inches(2.5), Inches(4.9), Inches(0.5),
             title, size=18, color=WHITE, bold=True)
    txBox = add_text(slide, Inches(x + 0.3), Inches(3.1), Inches(4.9), Inches(3.5), "", size=12)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for item in items:
        add_para(tf, f"▸  {item}", size=12, color=LIGHT_GRAY, space_before=Pt(6))


# ============================================================
# 슬라이드 6: Day 5~6 타임라인
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(0.8), ACCENT2)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
         "스프린트 타임라인 — Day 5 ~ 6", size=36, color=WHITE, bold=True)

days_3 = [
    ("Day 5", "D1 + R2 + 서비스", ACCENT2, [
        "D1 스키마: service_registry, news_cache, feature_flags",
        "시드 데이터: 서비스 2개, 뉴스 3개, 플래그 3개",
        "D1 접근 레이어 (services, news, flags)",
        "R2 헬퍼 (경로 빌더 + CRUD)",
        "/services 목록 (D1 기반, 상태 뱃지)",
        "/services/news 뉴스 애그리게이터 MVP",
    ]),
    ("Day 6", "홈 폴리싱", ACCENT3, [
        "2D 히어로: 도트 패턴 + 카피 + CTA",
        "키워드 태그 float 애니메이션 (CSS)",
        "프로젝트·Showcase 카드에 Link 추가",
        "서비스 섹션 추가 (D1 기반)",
        "구독 CTA (SubscribeForm 포함)",
        "광고 placeholder 컴포넌트",
    ]),
]

for i, (day, title, color, items) in enumerate(days_3):
    x = 0.8 + i * 6.0
    card = add_shape(slide, Inches(x), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD, 0.02)
    add_badge(slide, Inches(x + 0.3), Inches(2.0), day, color, BG_DARK, Inches(1.2))
    add_text(slide, Inches(x + 0.3), Inches(2.5), Inches(4.9), Inches(0.5),
             title, size=18, color=WHITE, bold=True)
    txBox = add_text(slide, Inches(x + 0.3), Inches(3.1), Inches(4.9), Inches(3.5), "", size=12)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for item in items:
        add_para(tf, f"▸  {item}", size=12, color=LIGHT_GRAY, space_before=Pt(6))


# ============================================================
# 슬라이드 7: Day 7
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(0.8), ACCENT4)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
         "스프린트 타임라인 — Day 7 (마감)", size=36, color=WHITE, bold=True)

card = add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8), BG_CARD, 0.02)
add_badge(slide, Inches(1.2), Inches(2.0), "Day 7", ACCENT4, BG_DARK, Inches(1.2))
add_text(slide, Inches(1.2), Inches(2.5), Inches(10), Inches(0.5),
         "SEO + QA + GitHub 정리 + 배포 준비", size=20, color=WHITE, bold=True)

col1_items = [
    ("SEO Metadata", "전역 metadataBase + title.template + OG/Twitter"),
    ("페이지별 Metadata", "정적 8개 + generateMetadata 3개"),
    ("sitemap.ts", "정적 6경로 + 동적 상세 (Velite 기반)"),
    ("robots.ts", "전체 허용, /api/ 차단"),
]
col2_items = [
    ("RSS 2.0", "/rss.xml Route Handler (blog 기반)"),
    ("README.md", "기술 스택, 실행 방법, 폴더 구조"),
    (".env.example", "SITE_URL + Supabase placeholder"),
    ("QA", "모바일 nav 수정 + ESLint ignore"),
]

for i, (title, desc) in enumerate(col1_items):
    y = 3.2 + i * 0.65
    add_text(slide, Inches(1.4), Inches(y), Inches(2), Inches(0.4),
             title, size=13, color=ACCENT, bold=True)
    add_text(slide, Inches(3.5), Inches(y), Inches(4), Inches(0.4),
             desc, size=12, color=LIGHT_GRAY)

for i, (title, desc) in enumerate(col2_items):
    y = 3.2 + i * 0.65
    add_text(slide, Inches(7.2), Inches(y), Inches(2), Inches(0.4),
             title, size=13, color=ACCENT2, bold=True)
    add_text(slide, Inches(9.3), Inches(y), Inches(4), Inches(0.4),
             desc, size=12, color=LIGHT_GRAY)


# ============================================================
# 슬라이드 8: 사이트 구조
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(0.8), ACCENT)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
         "사이트 구조 — 12개 라우트", size=36, color=WHITE, bold=True)

routes = [
    ("/", "홈", "Static", "히어로 + 블로그/프로젝트/Showcase/서비스/구독"),
    ("/blog", "블로그 목록", "Dynamic", "카테고리·태그 필터, searchParams"),
    ("/blog/[slug]", "블로그 상세", "SSG", "MDX 본문, TOC, 코드 하이라이팅"),
    ("/projects", "프로젝트 목록", "Static", "카드 목록 + 상세 링크"),
    ("/projects/[slug]", "프로젝트 상세", "SSG", "techStack, repo/demo 링크"),
    ("/showcase", "Showcase 목록", "Static", "상태 뱃지 + 카드"),
    ("/showcase/[slug]", "Showcase 상세", "SSG", "스택, 외부 링크"),
    ("/services", "서비스 목록", "Static", "D1 기반, 상태 뱃지"),
    ("/services/news", "뉴스", "Dynamic", "D1 news_cache, 카테고리 필터"),
    ("/about", "소개", "Static", "경력 타임라인, 기술 스택"),
    ("/subscribe", "구독", "Static", "SubscribeForm"),
]

# 테이블 헤더
headers = ["경로", "페이지", "렌더링", "주요 기능"]
header_widths = [2.2, 1.5, 1.2, 5.5]
x_start = 0.8
y_header = 1.75

for j, (h, w) in enumerate(zip(headers, header_widths)):
    x = x_start + sum(header_widths[:j])
    add_shape(slide, Inches(x), Inches(y_header), Inches(w), Inches(0.4), ACCENT, 0.01)
    add_text(slide, Inches(x + 0.1), Inches(y_header), Inches(w - 0.2), Inches(0.4),
             h, size=12, color=BG_DARK, bold=True, align=PP_ALIGN.LEFT)

for i, (route, name, render_type, desc) in enumerate(routes):
    y = y_header + 0.42 + i * 0.42
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x1F, 0x2B, 0x42)
    vals = [route, name, render_type, desc]
    for j, (val, w) in enumerate(zip(vals, header_widths)):
        x = x_start + sum(header_widths[:j])
        add_shape(slide, Inches(x), Inches(y), Inches(w), Inches(0.4), bg, 0.005)
        c = ACCENT if j == 0 else (ACCENT2 if j == 2 else LIGHT_GRAY)
        add_text(slide, Inches(x + 0.1), Inches(y), Inches(w - 0.2), Inches(0.4),
                 val, size=10, color=c, bold=(j == 0))


# ============================================================
# 슬라이드 9: 데이터 아키텍처
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(0.8), ACCENT3)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
         "데이터 아키텍처", size=36, color=WHITE, bold=True)

data_layers = [
    ("Velite (MDX)", ACCENT, "빌드 타임 컴파일", [
        "blog → 3개 글",
        "projects → 2개",
        "showcase → 3개",
        "draft 자동 필터링",
    ]),
    ("Cloudflare D1", ACCENT2, "런타임 SQLite", [
        "service_registry",
        "news_cache",
        "feature_flags",
        "try/catch fallback",
    ]),
    ("Cloudflare R2", ACCENT4, "오브젝트 스토리지", [
        "blog/{slug}/cover.ext",
        "projects/{slug}/thumb.ext",
        "og/{type}/{slug}.png",
        "헬퍼 구현 완료 (Phase 2 사용)",
    ]),
    ("Supabase", MID_GRAY, "Phase 2 예정", [
        "인증 (Auth)",
        "댓글 (Comments)",
        "구독 (Subscriptions)",
        "멤버십 (Memberships)",
    ]),
]

for i, (title, color, subtitle, items) in enumerate(data_layers):
    x = 0.8 + i * 3.1
    card = add_shape(slide, Inches(x), Inches(1.8), Inches(2.8), Inches(4.5), BG_CARD, 0.02)
    add_accent_line(slide, Inches(x + 0.3), Inches(2.0), Inches(0.6), color)
    add_text(slide, Inches(x + 0.3), Inches(2.2), Inches(2.2), Inches(0.4),
             title, size=16, color=color, bold=True)
    add_text(slide, Inches(x + 0.3), Inches(2.6), Inches(2.2), Inches(0.4),
             subtitle, size=11, color=MID_GRAY)
    txBox = add_text(slide, Inches(x + 0.3), Inches(3.1), Inches(2.2), Inches(2.8), "", size=12)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for item in items:
        add_para(tf, f"•  {item}", size=12, color=LIGHT_GRAY, space_before=Pt(6))


# ============================================================
# 슬라이드 10: AI 페어 프로그래밍 워크플로우
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(0.8), ACCENT3)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
         "AI 페어 프로그래밍 워크플로우", size=36, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
         "Claude Code를 활용한 체계적 개발 프로세스", size=16, color=MID_GRAY)

steps = [
    ("1", "문서 읽기", "CLAUDE.md\nplan.md\ntechspec.md", ACCENT),
    ("2", "계획 수립", "작업 범위 정의\n수정 파일 목록\n구현 전략", ACCENT2),
    ("3", "구현", "코드 생성\n린트 검증\n빌드 확인", ACCENT3),
    ("4", "검증", "pnpm lint\npnpm build\n수동 체크리스트", ACCENT4),
    ("5", "커밋", "상세 커밋 메시지\nGitHub push\nplan.md 갱신", ACCENT),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.5 + i * 2.5
    # 번호 원
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.6), Inches(2.3), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(24)
    circle.text_frame.paragraphs[0].font.color.rgb = BG_DARK
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 화살표 (마지막 제외)
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.5), Inches(2.5), Inches(0.7), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_GRAY
        arrow.line.fill.background()

    add_text(slide, Inches(x + 0.2), Inches(3.2), Inches(1.5), Inches(0.4),
             title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    card = add_shape(slide, Inches(x + 0.05), Inches(3.7), Inches(1.8), Inches(2.2), BG_CARD, 0.02)
    lines = desc.split("\n")
    txBox = add_text(slide, Inches(x + 0.2), Inches(3.9), Inches(1.5), Inches(1.8), "", size=11)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in lines:
        add_para(tf, line, size=11, color=LIGHT_GRAY, space_before=Pt(4))


# ============================================================
# 슬라이드 11: 향후 계획
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(0.8), ACCENT4)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7),
         "향후 계획 — Phase 2", size=36, color=WHITE, bold=True)

phase2 = [
    ("인증 · 사용자", ACCENT, [
        "Supabase Auth 연동",
        "댓글 시스템",
        "구독 백엔드 (MailerLite/Buttondown)",
        "멤버십 기능",
    ]),
    ("콘텐츠 확장", ACCENT2, [
        "뉴스 외부 수집 자동화 (크롤러/RSS)",
        "OG 이미지 동적 생성",
        "R2 이미지 업로드 파이프라인",
        "블로그 시리즈 기능",
    ]),
    ("UX 개선", ACCENT3, [
        "다크모드 토글",
        "3D 히어로 (성능 검증 후)",
        "모바일 햄버거 메뉴",
        "페이지 전환 애니메이션",
    ]),
    ("운영", ACCENT4, [
        "Cloudflare 프로덕션 배포",
        "커스텀 도메인 연결",
        "Lighthouse 성능 최적화",
        "GitHub Actions CI/CD",
    ]),
]

for i, (title, color, items) in enumerate(phase2):
    x = 0.8 + i * 3.1
    card = add_shape(slide, Inches(x), Inches(1.8), Inches(2.8), Inches(4.2), BG_CARD, 0.02)
    add_accent_line(slide, Inches(x + 0.3), Inches(2.0), Inches(0.6), color)
    add_text(slide, Inches(x + 0.3), Inches(2.2), Inches(2.2), Inches(0.5),
             title, size=16, color=color, bold=True)
    txBox = add_text(slide, Inches(x + 0.3), Inches(2.8), Inches(2.2), Inches(3), "", size=12)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for item in items:
        add_para(tf, f"▸  {item}", size=12, color=LIGHT_GRAY, space_before=Pt(6))


# ============================================================
# 슬라이드 12: 마무리
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(-1), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
circle.line.fill.background()

add_accent_line(slide, Inches(4.5), Inches(2.5), Inches(1.2), ACCENT)
add_text(slide, Inches(2), Inches(2.8), Inches(9), Inches(1),
         "만들고, 기록하고, 공유합니다.", size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.6),
         "gotech.lab — 1주 MVP 스프린트 완료", size=20, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(4.6), Inches(9), Inches(0.5),
         "github.com/kyeongseok-go/dev-gotech-lab", size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)

# 하단 뱃지들
badges_data = ["Day 0~7 완료", "12개 라우트", "8개 MDX", "D1 + R2", "SEO + RSS"]
badge_x_start = 2.5
for i, badge_text in enumerate(badges_data):
    colors = [ACCENT, ACCENT2, ACCENT3, ACCENT4, ACCENT]
    add_badge(slide, Inches(badge_x_start + i * 1.8), Inches(5.5), badge_text, colors[i], BG_DARK, Inches(1.5))


# ============================================================
# 저장
# ============================================================
output_path = "/Users/eunsuk.ko128/workspace/dev-gotech-lab/docs/output/gotech-lab-sprint-report.pptx"
prs.save(output_path)
print(f"PPT 저장 완료: {output_path}")
